from pathlib import Path

import pandas as pd
import pymupdf as fitz
from docx import Document
from loguru import logger
from pptx import Presentation  

from app.ingestion.exceptions import UnsupportedFileType
from app.ingestion.schemas import ParsedDocument, ParsedPage


class DocumentParser:

    def __init__(self):
        pass

    def parse(self, file_path: str | Path) -> ParsedDocument:

        self.file_path = Path(file_path)
        self.extension = self.file_path.suffix.lower()

        logger.info(f"Parsing document: {self.file_path.name}")

        if self.extension == ".pdf":
            return self._parse_pdf()

        elif self.extension == ".docx":
            return self._parse_docx()

        # elif self.extension == ".pptx":  # DISABLED: pptx module not available
        #     return self._parse_pptx()

        elif self.extension in [".xlsx", ".xls"]:
            return self._parse_xlsx()

        elif self.extension == ".csv":
            return self._parse_csv()

        elif self.extension in [".txt", ".md", ".html", ".htm"]:
            return self._parse_txt()

        raise UnsupportedFileType(
            f"{self.extension} is not supported."
        )

    def _parse_pdf(self) -> ParsedDocument:

        try:
            parsed_pages = []

            with fitz.open(self.file_path) as doc:
                for page_number, page in enumerate(doc, start=1):
                    text = page.get_text().strip()
                    parsed_pages.append(
                        ParsedPage(
                            document=self.file_path.name,
                            page=page_number,
                            text=text,
                            needs_ocr=(
                                len(text) < 50
                                or len(text.split()) < 20
                            ),
                        )
                    )

            document = ParsedDocument(
                document=self.file_path.name,
                file_type="pdf",
                page_count=len(parsed_pages),
                total_characters=sum(
                    len(page.text)
                    for page in parsed_pages
                ),
                needs_ocr=any(
                    page.needs_ocr
                    for page in parsed_pages
                ),
                ocr_used=False,
                pages=parsed_pages,
            )

            logger.success(
                f"PDF parsed successfully ({len(parsed_pages)} pages)"
            )

            return document

        except Exception as e:
            logger.exception(e)
            raise RuntimeError(
                f"Failed to parse PDF: {e}"
            )

    def _parse_docx(self) -> ParsedDocument:

        try:
            document = Document(self.file_path)
            text = "\n".join(
                paragraph.text.strip()
                for paragraph in document.paragraphs
                if paragraph.text.strip()
            )

            pages = [
                ParsedPage(
                    document=self.file_path.name,
                    page=1,
                    text=text,
                    needs_ocr=False,
                )
            ]

            return ParsedDocument(
                document=self.file_path.name,
                file_type="docx",
                page_count=1,
                total_characters=len(text),
                needs_ocr=False,
                ocr_used=False,
                pages=pages,
            )

        except Exception as e:
            logger.exception(e)
            raise RuntimeError(
                f"Failed to parse DOCX: {e}"
            )

    def _parse_pptx(self) -> ParsedDocument:

        try:
            presentation = Presentation(self.file_path)
            slides = []

            for index, slide in enumerate(
                presentation.slides,
                start=1,
            ):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())

                slide_text = "\n".join(texts)

                slides.append(
                    ParsedPage(
                        document=self.file_path.name,
                        page=index,
                        text=slide_text,
                        needs_ocr=False,
                    )
                )

            return ParsedDocument(
                document=self.file_path.name,
                file_type="pptx",
                page_count=len(slides),
                total_characters=sum(
                    len(slide.text)
                    for slide in slides
                ),
                needs_ocr=False,
                ocr_used=False,
                pages=slides,
            )

        except Exception as e:
            logger.exception(e)
            raise RuntimeError(
                f"Failed to parse PPTX: {e}"
            )

    def _parse_xlsx(self) -> ParsedDocument:

        try:
            workbook = pd.read_excel(
                self.file_path,
                sheet_name=None,
            )

            sheets = []

            for sheet_name, dataframe in workbook.items():
                sheet_text = dataframe.to_string(index=False)

                sheets.append(
                    ParsedPage(
                        document=self.file_path.name,
                        page=str(sheet_name),
                        text=sheet_text,
                        needs_ocr=False,
                    )
                )

            return ParsedDocument(
                document=self.file_path.name,
                file_type="xlsx",
                page_count=len(sheets),
                total_characters=sum(
                    len(sheet.text)
                    for sheet in sheets
                ),
                needs_ocr=False,
                ocr_used=False,
                pages=sheets,
            )

        except Exception as e:
            logger.exception(e)
            raise RuntimeError(
                f"Failed to parse Excel spreadsheet: {e}"
            )

    def _parse_csv(self) -> ParsedDocument:

        try:
            dataframe = pd.read_csv(self.file_path)
            
            # Create separate pages for each row to enable proper chunking
            # This ensures each row is treated as independent data
            pages = []
            
            # Add header row as first page
            header_text = ", ".join(dataframe.columns.tolist())
            pages.append(
                ParsedPage(
                    document=self.file_path.name,
                    page=1,
                    text=header_text,
                    needs_ocr=False,
                )
            )
            
            # Add each data row as a separate page
            for row_idx, (_, row) in enumerate(dataframe.iterrows(), start=2):
                # Convert row to readable string format
                row_text = ", ".join(
                    f"{col}: {str(val).strip()}"
                    for col, val in row.items()
                    if pd.notna(val)
                )
                pages.append(
                    ParsedPage(
                        document=self.file_path.name,
                        page=row_idx,
                        text=row_text,
                        needs_ocr=False,
                    )
                )
            
            total_text = "\n".join(page.text for page in pages)

            return ParsedDocument(
                document=self.file_path.name,
                file_type="csv",
                page_count=len(pages),
                total_characters=len(total_text),
                needs_ocr=False,
                ocr_used=False,
                pages=pages,
            )

        except Exception as e:
            logger.exception(e)
            raise RuntimeError(
                f"Failed to parse CSV document: {e}"
            )

    def _parse_txt(self) -> ParsedDocument:

        try:
            text = self.file_path.read_text(encoding="utf-8", errors="ignore")

            pages = [
                ParsedPage(
                    document=self.file_path.name,
                    page=1,
                    text=text,
                    needs_ocr=False,
                )
            ]

            return ParsedDocument(
                document=self.file_path.name,
                file_type=self.extension.lstrip("."),
                page_count=1,
                total_characters=len(text),
                needs_ocr=False,
                ocr_used=False,
                pages=pages,
            )

        except Exception as e:
            logger.exception(e)
            raise RuntimeError(
                f"Failed to parse text document: {e}"
            )